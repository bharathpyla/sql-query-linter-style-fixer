import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Modern Professional Indigo & Violet Light Theme Palette
COLOR_BG = RGBColor(230, 237, 246)       # Soft Sky/Ice Blue canvas background (dullness fix)
COLOR_TITLE = RGBColor(15, 23, 42)       # Slate 900 (High-contrast deep slate for main titles)
COLOR_BODY = RGBColor(51, 65, 85)        # Slate 700 (Clear readable charcoal body text)
COLOR_ACCENT = RGBColor(79, 70, 229)     # Indigo 600 (Vibrant Royal Indigo for highlights/dividers)
COLOR_HIGHLIGHT = RGBColor(219, 39, 119) # Pink 600 (Eye-catching fuchsia for team header)
COLOR_MUTED = RGBColor(100, 116, 139)    # Slate 500 (Muted gray for footers)
COLOR_CARD_BG = RGBColor(255, 255, 255)  # Pure White background for card layouts
COLOR_CARD_BORDER = RGBColor(224, 231, 255) # Indigo 100 (Very soft card border color)

COLOR_PROBLEM_ACCENT = RGBColor(239, 68, 68) # Red 500 (Vibrant red accent line for problems)
COLOR_SUCCESS_ACCENT = RGBColor(16, 185, 129) # Emerald 500 (Vibrant green accent line for advantages/conclusions)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_divider_line(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # Remove shape border
    return shape

def add_footer(slide):
    footer_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.95), Inches(11.83), Inches(0.3))
    tf = footer_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "SQL Query Linter & Style Fixer  |  Developed by Team 20"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.LEFT

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    
    # Elegant thin royal indigo divider line under title
    add_divider_line(slide, Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.04), COLOR_ACCENT)
    
    # Add slide footer
    add_footer(slide)
    return title_box

def add_text_card(slide, left, top, width, height, title, desc, accent_color=COLOR_ACCENT):
    # Rounded rectangle background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    shape.line.color.rgb = COLOR_CARD_BORDER
    shape.line.width = Pt(1.5)
    
    # Left vertical accent line on the card
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Text box
    tb = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(0.08), width - Inches(0.55), height - Inches(0.16))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"{title}: "
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    
    run = p.add_run()
    run.text = desc
    run.font.name = 'Segoe UI'
    run.font.size = Pt(13)
    run.font.bold = False
    run.font.color.rgb = COLOR_BODY

def add_step_card(slide, left, top, width, height, step_num, text, accent_color=COLOR_ACCENT):
    # Rounded rectangle background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    shape.line.color.rgb = COLOR_CARD_BORDER
    shape.line.width = Pt(1.5)
    
    # Left vertical accent line on the card
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Step badge text box
    badge_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.11), Inches(0.8), height - Inches(0.22))
    tf_badge = badge_box.text_frame
    tf_badge.word_wrap = True
    tf_badge.margin_left = tf_badge.margin_top = tf_badge.margin_right = tf_badge.margin_bottom = 0
    p_badge = tf_badge.paragraphs[0]
    p_badge.text = step_num
    p_badge.font.name = 'Segoe UI'
    p_badge.font.size = Pt(17)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_ACCENT
    
    # Description text box
    tb = slide.shapes.add_textbox(left + Inches(1.0), top + Inches(0.11), width - Inches(1.15), height - Inches(0.22))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_BODY

def add_grid_card(slide, left, top, width, height, title, desc, accent_color=COLOR_ACCENT):
    # Rounded rectangle background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    shape.line.color.rgb = COLOR_CARD_BORDER
    shape.line.width = Pt(1.5)
    
    # Left vertical accent line on the card
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Text box
    tb = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(0.2), width - Inches(0.55), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_BODY

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Premium Indigo landing page)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)
    
    # Elegant top brand bar
    add_divider_line(slide1, Inches(0), Inches(0), Inches(13.33), Inches(0.12), COLOR_ACCENT)
    
    # Large Title Box
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.83), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    p1.text = "SQL Query Linter & Style Fixer"
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    
    p2 = tf.add_paragraph()
    p2.text = "Hybrid Local-AI Database Optimization Pipeline"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    p2.space_before = Pt(10)
    
    # Team Info Box (Nested inside a gorgeous clean card container)
    team_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.9), Inches(11.83), Inches(2.7))
    team_card.fill.solid()
    team_card.fill.fore_color.rgb = COLOR_CARD_BG
    team_card.line.color.rgb = COLOR_CARD_BORDER
    team_card.line.width = Pt(1.5)
    
    # Left border highlight on the team card
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.9), Inches(0.12), Inches(2.7))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_HIGHLIGHT
    accent_bar.line.fill.background()
    
    team_box = slide1.shapes.add_textbox(Inches(1.15), Inches(4.1), Inches(11.23), Inches(2.3))
    tf_team = team_box.text_frame
    tf_team.word_wrap = True
    tf_team.margin_left = tf_team.margin_top = tf_team.margin_right = tf_team.margin_bottom = 0
    
    p_team_hdr = tf_team.paragraphs[0]
    p_team_hdr.text = "DEVELOPED BY: TEAM 20"
    p_team_hdr.font.name = 'Segoe UI'
    p_team_hdr.font.size = Pt(13)
    p_team_hdr.font.bold = True
    p_team_hdr.font.color.rgb = COLOR_HIGHLIGHT
    p_team_hdr.space_after = Pt(6)
    
    members = [
        "Penta Satwika (23U41A4245) - Lead Streamlit UI & Dashboard Developer",
        "Bharath Pyla (23U41A0536) - Core SQL Parser & Spacing Engine Architect",
        "Aythi Eswar (23U41A0406) - Gemini LLM Prompt Engineer & AI Architect",
        "Dadi James (24U45A0212) - AST Validation Loop & QA Tester"
    ]
    
    for member in members:
        p_mem = tf_team.add_paragraph()
        p_mem.text = member
        p_mem.font.name = 'Segoe UI'
        p_mem.font.size = Pt(13)
        p_mem.font.color.rgb = COLOR_BODY
        p_mem.space_before = Pt(3)
        
    # ----------------------------------------------------
    # SLIDE 2: Objective (Beautiful Two-Column Layout)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_title(slide2, "Objective")
    
    # Left Box - Large Styled Summary Card
    left_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.2), Inches(4.7))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(215, 226, 250)  # Rich soft Indigo/Blue tint
    left_card.line.color.rgb = COLOR_CARD_BORDER
    left_card.line.width = Pt(1.5)
    
    left_accent = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(0.12), Inches(4.7))
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = COLOR_ACCENT
    left_accent.line.fill.background()
    
    summary_box = slide2.shapes.add_textbox(Inches(1.15), Inches(2.2), Inches(4.5), Inches(3.9))
    tf_summary = summary_box.text_frame
    tf_summary.word_wrap = True
    tf_summary.margin_left = tf_summary.margin_top = tf_summary.margin_right = tf_summary.margin_bottom = 0
    p_sum = tf_summary.paragraphs[0]
    p_sum.text = "To build a robust database utility that automates layout formatting, syntax cleansing, and structural optimization for database development teams by merging client-side deterministic parsers with serverless Large Language Models."
    p_sum.font.name = 'Segoe UI'
    p_sum.font.size = Pt(19)
    p_sum.font.bold = True
    p_sum.font.color.rgb = COLOR_ACCENT
    p_sum.line_spacing = 1.3
    
    # Right Box - Objectives Bullet points
    right_box = slide2.shapes.add_textbox(Inches(6.45), Inches(1.8), Inches(6.13), Inches(4.7))
    tf2 = right_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    
    objectives = [
        "Standardize mixed casings (camelCase/PascalCase) into uniform lowercase snake_case.",
        "Enforce SQL keywords capitalization rules locally at zero compute cost.",
        "Isolate and purge markdown formatting artifacts to safeguard downstream compiler integrity.",
        "Refactor implicit comma-joins and nested subqueries into explicit joins and Common Table Expressions (CTEs) via AI logic.",
        "Secure syntax correctness using an in-memory Abstract Syntax Tree validation loop."
    ]
    
    for i, obj in enumerate(objectives):
        p_bullet = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p_bullet.text = f"\u2022  {obj}"
        p_bullet.font.name = 'Segoe UI'
        p_bullet.font.size = Pt(15)
        p_bullet.font.color.rgb = COLOR_BODY
        p_bullet.space_before = Pt(10)
        p_bullet.line_spacing = 1.2
        
    # ----------------------------------------------------
    # SLIDE 3: Problem Statement (Red-Accented Card List)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_title(slide3, "Problem Statement")
    
    problems = [
        ("Inconsistent Casing", "Mixed naming structures (camelCase, PascalCase, lowercase) create style violations and audit complications."),
        ("Legacy Join Syntax", "Extensive usage of implicit comma-joins (e.g., FROM TableA, TableB WHERE A.id = B.id) obscures join pathways and increases maintenance overhead."),
        ("Complex Subquery Nesting", "Deeply nested SELECT projections decrease readability and limit database optimizer index lookup performance."),
        ("Cost Inflation", "Unrestricted use of SELECT * drives up cloud data processing costs under pay-per-query models."),
        ("Validation Vulnerability", "Rule-based linters only clean basic whitespace, whereas raw LLM optimization agents can hallucinate and generate syntactically broken SQL queries.")
    ]
    
    card_y = 1.8
    card_h = 0.85
    card_gap = 0.1
    for prob_title, prob_desc in problems:
        add_text_card(slide3, Inches(0.75), Inches(card_y), Inches(11.83), Inches(card_h), prob_title, prob_desc, COLOR_PROBLEM_ACCENT)
        card_y += card_h + card_gap
        
    # ----------------------------------------------------
    # SLIDE 4: Technologies Used (Indigo-Accented Card List)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_title(slide4, "Technologies Used")
    
    techs = [
        ("Python", "Core language for deterministic regex compilation, script execution, and API routing."),
        ("Streamlit Framework", "For designing the user interface dashboard workspace, the batch file uploader, and in-memory downloads."),
        ("Google Gemini API (gemini-2.5-flash)", "Serverless AI model utilized for semantic query refactoring (CTEs, Joins) configured with low temperature (0.1) and few-shot examples."),
        ("sqlglot Parser", "Local dialect-aware parser used to pretty-print statements and validate AST syntax inside the correction loop."),
        ("difflib", "Native engine used to compile unified diff reports between raw query inputs and optimized outputs.")
    ]
    
    card_y = 1.8
    card_h = 0.85
    card_gap = 0.1
    for tech, desc in techs:
        add_text_card(slide4, Inches(0.75), Inches(card_y), Inches(11.83), Inches(card_h), tech, desc, COLOR_ACCENT)
        card_y += card_h + card_gap
        
    # ----------------------------------------------------
    # SLIDE 5: System Architecture & Workflow (Sequential Roadmap)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_title(slide5, "System Architecture & Workflow")
    
    flow_steps = [
        ("01", "Query Input -> Sandbox Area, File Uploader, or CLI batch scanner reads statement."),
        ("02", "Input Sanitization -> strip_markdown_artifacts() strips code blocks, bold asterisks, and italic underscores safely."),
        ("03", "Phase 1 Linter -> Deterministic capitalization of standard keywords, regex-based conversion of camelCase to snake_case, and SELECT * diagnostic checks."),
        ("04", "Phase 2 Refactor -> gemini-2.5-flash receives baseline formatting and optimizes syntax structures (implicit joins to explicit joins, subqueries to CTEs)."),
        ("05", "Self-Correction Loop -> Output SQL is evaluated by local parser. If syntax errors occur, the traceback is fed back to the LLM to self-heal (up to 3 iterations)."),
        ("06", "Finalization -> Renders refactored SQL inside copyable st.code frames and outputs colored diff blocks.")
    ]
    
    card_y = 1.8
    card_h = 0.70
    card_gap = 0.1
    for step_num, text in flow_steps:
        add_step_card(slide5, Inches(0.75), Inches(card_y), Inches(11.83), Inches(card_h), step_num, text, COLOR_ACCENT)
        card_y += card_h + card_gap
        
    # ----------------------------------------------------
    # SLIDE 6: Modules Overview (Vibrant Vertical Card Layout)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6)
    add_title(slide6, "Modules Overview")
    
    modules = [
        ("Sanitization & Cleaning Module", "Contains lookbehind/lookahead regular expressions that clean styling artifacts from incoming text without corrupting snake_case columns."),
        ("Phase 1: Local Linter Module", "Zero-cost deterministic casing cleaner and capitalization compiler. Enforces standard keywords and casing conventions locally before network calls."),
        ("Phase 2: AI Refactoring Module", "Leverages Gemini LLM with few-shot instructions to translate legacy comma-joins and extract subqueries into CTEs."),
        ("Validation & Critic Module", "Compares LLM outputs against local syntax rules using sqlglot, orchestrating the self-correcting logic loop if discrepancies are found."),
        ("Interface & Download Module", "Manages dual Streamlit and CLI interfaces, in-memory uploads, diff blocks rendering, and download actions.")
    ]
    
    card_y = 1.8
    card_h = 0.85
    card_gap = 0.1
    for mod_title, mod_desc in modules:
        add_text_card(slide6, Inches(0.75), Inches(card_y), Inches(11.83), Inches(card_h), mod_title, mod_desc, COLOR_ACCENT)
        card_y += card_h + card_gap
        
    # ----------------------------------------------------
    # SLIDE 7: Project Advantages (Dynamic 2x2 Grid Card Layout)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7)
    add_title(slide7, "Project Advantages")
    
    advs = [
        ("Cost-Efficiency", "Capitalization and casing rules run locally at zero token cost, conserving expensive LLM API bandwidth for semantic restructuring tasks."),
        ("Syntax Security", "The self-correcting validation loop acts as a compiler guard, preventing code hallucinations and ensuring the user only receives valid, executable SQL."),
        ("Cloud Deployability", "Supports in-memory batch file uploads, allowing the application to run smoothly on isolated cloud deployment servers like Streamlit Community Cloud."),
        ("Pipeline Integration", "Dual CLI and web dashboard structure allows the tool to fit seamlessly into single-query debugging sessions or automated Git pre-commit hooks.")
    ]
    
    # 2x2 grid positioning
    add_grid_card(slide7, Inches(0.75), Inches(1.8), Inches(5.7), Inches(2.2), advs[0][0], advs[0][1], COLOR_SUCCESS_ACCENT)
    add_grid_card(slide7, Inches(6.88), Inches(1.8), Inches(5.7), Inches(2.2), advs[1][0], advs[1][1], COLOR_SUCCESS_ACCENT)
    add_grid_card(slide7, Inches(0.75), Inches(4.2), Inches(5.7), Inches(2.2), advs[2][0], advs[2][1], COLOR_SUCCESS_ACCENT)
    add_grid_card(slide7, Inches(6.88), Inches(4.2), Inches(5.7), Inches(2.2), advs[3][0], advs[3][1], COLOR_SUCCESS_ACCENT)
    
    # ----------------------------------------------------
    # SLIDE 8: Conclusion (Beautiful Two-Column Layout)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8)
    add_title(slide8, "Conclusion")
    
    # Left Box - Large Styled Conclusion Statement Card
    left_card8 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(5.2), Inches(4.7))
    left_card8.fill.solid()
    left_card8.fill.fore_color.rgb = RGBColor(209, 250, 229)  # Rich soft Emerald/Green tint
    left_card8.line.color.rgb = COLOR_CARD_BORDER
    left_card8.line.width = Pt(1.5)
    
    left_accent8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(0.12), Inches(4.7))
    left_accent8.fill.solid()
    left_accent8.fill.fore_color.rgb = COLOR_SUCCESS_ACCENT
    left_accent8.line.fill.background()
    
    concl_box = slide8.shapes.add_textbox(Inches(1.15), Inches(2.2), Inches(4.5), Inches(3.9))
    tf_concl = concl_box.text_frame
    tf_concl.word_wrap = True
    tf_concl.margin_left = tf_concl.margin_top = tf_concl.margin_right = tf_concl.margin_bottom = 0
    p_concl = tf_concl.paragraphs[0]
    p_concl.text = "The SQL Query Linter & Style Fixer successfully bridges the gap between rule-based style cleaners and AI-driven semantic query model refactoring."
    p_concl.font.name = 'Segoe UI'
    p_concl.font.size = Pt(20)
    p_concl.font.bold = True
    p_concl.font.color.rgb = COLOR_SUCCESS_ACCENT
    p_concl.line_spacing = 1.3
    
    # Right Box - Conclusion points list
    right_box8 = slide8.shapes.add_textbox(Inches(6.45), Inches(1.8), Inches(6.13), Inches(4.7))
    tf8 = right_box8.text_frame
    tf8.word_wrap = True
    tf8.margin_left = tf8.margin_top = tf8.margin_right = tf8.margin_bottom = 0
    
    points = [
        "Optimizes developer workflow speed and audit readiness for legacy SQL codebases.",
        "Guarantees compilation integrity through an AST validation loop architecture.",
        "Reduces database runtime search complexity by converting nested subqueries to CTEs and implicit comma-joins to explicit joins.",
        "Provides a cloud-deployable dashboard and automated CLI script to accommodate any integration pipeline environment."
    ]
    
    for i, pt in enumerate(points):
        p_pt = tf8.add_paragraph() if i > 0 else tf8.paragraphs[0]
        p_pt.text = f"\u2022  {pt}"
        p_pt.font.name = 'Segoe UI'
        p_pt.font.size = Pt(15)
        p_pt.font.color.rgb = COLOR_BODY
        p_pt.space_before = Pt(10)
        p_pt.line_spacing = 1.2
        
    prs.save("Presentation.pptx")
    print("Presentation saved successfully.")

if __name__ == '__main__':
    build_presentation()
